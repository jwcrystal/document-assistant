import os
import tempfile
import fitz  # PyMuPDF
import unicodedata
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
import re
from nicegui import ui, app
from pathlib import Path
from docling.document_converter import DocumentConverter
import logging
import asyncio
from typing import Optional, Callable
import pandas as pd

# 設定 logging
logging.basicConfig(level=logging.INFO)
logging.getLogger('docling').setLevel(logging.WARNING)
logger = logging.getLogger(__name__)

# 創建臨時上傳目錄
UPLOAD_DIR = Path("temp_uploads")
UPLOAD_DIR.mkdir(exist_ok=True)

# 全局變數
MAX_FILE_SIZE = 200_000_000  # 200MB
current_file_path = None
pdf_pages = 0
current_page = 0
ocr_result = None
loading = None  # 全局加載狀態
preview_container = None  # 預覽容器
ocr_result_container = None  # OCR 結果容器

def init_application():
    """初始化應用程式狀態"""
    global current_file_path, pdf_pages, current_page, ocr_result, loading, preview_container, ocr_result_container
    
    # 重置全局變數
    current_file_path = None
    pdf_pages = 0
    current_page = 0
    ocr_result = None
    loading = None
    
    # 清除容器內容
    if preview_container is not None:
        preview_container.clear()
    
    if ocr_result_container is not None:
        ocr_result_container.clear()
    
    try:
        # 清空上傳目錄
        for file in UPLOAD_DIR.glob('*'):
            try:
                if file.is_file():
                    file.unlink()
            except Exception as e:
                logger.error(f"無法刪除文件 {file}: {e}")
        
        logger.debug("應用程式狀態已重置")
    except Exception as e:
        logger.error(f"初始化應用程式時出錯: {e}")

# 檔名標準化函數
def sanitize_filename(filename: str) -> str:
    """將檔名標準化，移除特殊字元和空格"""
    # 使用 NFKD 標準化並編碼為 UTF-8，忽略無法處理的字元
    normalized = unicodedata.normalize('NFKD', filename).encode('utf-8', 'ignore').decode('utf-8')
    sanitized = normalized.replace(" ", "_")
    # 移除所有非字母、數字、底線、連字號和點的字元
    sanitized = re.sub(r'[^\w.-]', '', sanitized)
    return sanitized

async def handle_upload(e):
    """處理文件上傳"""
    global current_file_path, pdf_pages, current_page, preview_container
    
    # 重置狀態
    current_page = 0
    if preview_container is not None:
        preview_container.clear()
    
    # 檢查是否有文件上傳
    if not e.content:
        ui.notify("未選擇文件", type='warning')
        return
    
    try:
        # 檢查檔案大小
        file_size = len(e.content.read())
        e.content.seek(0)  # 重置文件指針
        
        if file_size > MAX_FILE_SIZE:
            ui.notify(f"檔案大小超過限制 (最大 {MAX_FILE_SIZE/1_000_000}MB)", type='negative')
            # 關閉並刪除暫存檔案
            e.content.close()
            # 清除上傳狀態
            if hasattr(e.sender, 'reset'):
                e.sender.reset()
            return
            
        # 獲取上傳的文件
        file_obj = e.content  # 獲取 SpooledTemporaryFile 物件
        file_name = e.name
        file_type = e.type
        
        # 使用臨時文件保存上傳的內容
        with tempfile.NamedTemporaryFile(
            delete=False, 
            suffix=os.path.splitext(file_name)[1]  # 保留原始副檔名
        ) as tmp:
            # 重置文件指針並寫入文件內容
            file_obj.seek(0)
            tmp.write(file_obj.read())
            tmp_path = tmp.name  # 獲取臨時文件路徑
        
        try:
            # 獲取文件大小
            file_size = os.path.getsize(tmp_path)
            
            # 生成安全檔名並移動到目標目錄
            safe_name = sanitize_filename(file_name)
            file_path = UPLOAD_DIR / safe_name
            os.rename(tmp_path, file_path)  # 移動文件到目標位置
            
            # 更新全局變量
            current_file_path = file_path
            pdf_pages = 0
            current_page = 0
            
        except Exception as ex:
            # 發生錯誤時確保刪除臨時文件
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)
            raise ex
        
        # 顯示文件訊息
        file_info_text = f"""
        ### 文件資訊
        - 原始檔名: {file_name}
        - 安全檔名: {safe_name}
        - 文件大小: {file_size:,} bytes ({file_size / 1024 / 1024:.2f} MB)
        - 文件類型: {file_type}
        """
        
        if preview_container is None:
            preview_container = ui.column().classes('w-full q-mt-lg')
        
        with preview_container:
            ui.markdown(file_info_text)
            
            # 根據文件類型顯示預覽
            if file_type.startswith('image/'):
                with ui.image(str(file_path)).classes('w-full'):
                    pass
            elif file_type == 'application/pdf':
                await show_pdf_preview(file_path)
            elif file_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'application/msword']:
                await show_docx_preview(file_path)
            elif file_type == 'text/markdown' or file_name.lower().endswith(('.md', '.markdown')):
                await show_markdown_preview(file_path)
            elif file_type == 'text/html' or file_name.lower().endswith(('.htm', '.html')):
                await show_html_preview(file_path)
            elif file_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation', 'application/vnd.ms-powerpoint']:
                await show_pptx_preview(file_path)
            elif file_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', 'application/vnd.ms-excel']:
                await show_xlsx_preview(file_path)
            elif file_type == 'text/csv' or file_name.lower().endswith('.csv'):
                await show_csv_preview(file_path)
            else:
                ui.notify(f"不支援預覽 {file_name} 格式的文件", type='info')
            
        # 添加 OCR 按鈕
        with preview_container:
            with ui.row().classes('w-full justify-center mt-4'):
                ui.button('執行 OCR 辨識', on_click=lambda: run_ocr(file_path, safe_name), 
                         icon='image_search').props('color=primary')
                
    except Exception as ex:
        ui.notify(f"處理文件時發生錯誤: {str(ex)}", type='negative')

async def process_ocr(file_path: Path, progress_callback: Optional[Callable[[int, str], None]] = None):
    """執行 OCR 處理"""
    try:
        if progress_callback:
            await asyncio.get_event_loop().run_in_executor(None, lambda: progress_callback(10, "正在初始化..."))
            await asyncio.sleep(0.1)  # 讓 UI 有時間更新
        
        # 初始化轉換器
        converter = DocumentConverter()
        
        if progress_callback:
            await asyncio.get_event_loop().run_in_executor(None, lambda: progress_callback(30, "正在處理文件..."))
            await asyncio.sleep(0.1)
        
        # 執行轉換
        result = await asyncio.get_event_loop().run_in_executor(None, lambda: converter.convert(str(file_path)))

        if progress_callback:
            await asyncio.get_event_loop().run_in_executor(None, lambda: progress_callback(80, "正在生成 Markdown..."))
            await asyncio.sleep(0.1)
        
        # 轉換為 Markdown
        markdown_result = await asyncio.get_event_loop().run_in_executor(
            None, lambda: result.document.export_to_markdown()
        )
        
        if progress_callback:
            await asyncio.get_event_loop().run_in_executor(None, lambda: progress_callback(100, "處理完成!"))
            await asyncio.sleep(0.1)
        
        return markdown_result
    except UnicodeDecodeError as ude:
        error_msg = f"處理出錯: 文件格式不支援或已損壞"
        if progress_callback:
            await asyncio.get_event_loop().run_in_executor(None, lambda: progress_callback(0, error_msg))
        raise ValueError(error_msg) from ude
    except asyncio.CancelledError:
        raise
    except Exception as e:
        if progress_callback:
            await asyncio.get_event_loop().run_in_executor(None, lambda: progress_callback(0, f"處理出錯: {str(e)}"))
            await asyncio.sleep(0.1)
        raise
    finally:
        # 確保臨時文件清理
        try:
            if os.path.exists(file_path):
                os.unlink(file_path)
                logger.debug(f"已刪除臨時文件: {file_path}")
        except Exception as e:
            logger.error(f"刪除臨時文件 {file_path} 時出錯: {str(e)}")

# 定義全局變量用於保存 OCR 結果
ocr_result = None

# 全局變量用於控制 OCR 處理
ocr_task = None
ocr_cancelled = False

async def run_ocr(file_path: Path, original_filename: str):
    """執行 OCR 處理並顯示結果"""
    global ocr_task, ocr_cancelled
    
    # 重置取消標記
    ocr_cancelled = False
    
    # 創建進度對話框
    with ui.dialog() as dialog, ui.card().classes('w-96'):
        # 標題
        ui.label('OCR 處理中...').classes('text-lg font-bold')
        
        # 進度條
        progress = ui.linear_progress(0).props('instant-false')
        
        # 狀態文字
        status = ui.label('正在準備...')
        
        # 進度百分比
        percent = ui.label('0%')
        
        # 取消按鈕
        cancel_btn = ui.button('取消', on_click=lambda: cancel_ocr(dialog))
        cancel_btn.classes('mt-4')
        
        # 更新進度的回調函數
        def update_progress(value: int, message: str):
            progress.value = value / 100
            percent.text = f"{value}%"
            status.text = message
            ui.update(progress)
            ui.update(percent)
            ui.update(status)
        
        # 顯示對話框
        dialog.open()
        
        try:
            # 在背景執行 OCR 處理
            ocr_task = asyncio.create_task(process_ocr(file_path, update_progress))
            
            # 等待 OCR 處理完成或取消
            try:
                ocr_result = await ocr_task
                if ocr_cancelled:
                    return
                    
                # 關閉進度對話框
                dialog.close()
                
                # 顯示完成通知
                ui.notify("OCR 處理完成！", type='positive')
                
                show_ocr_result(ocr_result, original_filename)
                # 滾動到頁面底部
                ui.run_javascript('window.scrollTo({top: document.body.scrollHeight, behavior: "smooth"})')
                
            except asyncio.CancelledError:
                ui.notify("已取消 OCR 處理", type='warning')
            except Exception as e:
                if not ocr_cancelled:  # 只有當不是用戶取消時才顯示錯誤
                    dialog.close()
                    ui.notify(f"OCR 處理失敗: {str(e)}", type='negative')
            
        except Exception as e:
            dialog.close()
            ui.notify(f"執行 OCR 時發生錯誤: {str(e)}", type='negative')

def cancel_ocr(dialog):
    """取消正在進行的 OCR 處理"""
    global ocr_task, ocr_cancelled
    
    if ocr_task and not ocr_task.done():
        ocr_cancelled = True
        ocr_task.cancel()
        ui.notify("正在取消 OCR 處理...", type='warning')
    
    dialog.close()

def show_ocr_result(content: str, original_filename: str):
    """顯示 OCR 處理結果"""
    logger.debug(f"[DEBUG] 顯示 OCR 結果，內容類型: {type(content)}")
    if not content:
        ui.notify("處理錯誤: 無法辨識，請檢查檔案內容", type='negative')
        return
    
    # 確保結果容器存在
    global ocr_result_container
    if ocr_result_container is None:
        # 在 preview_container 下方創建結果容器
        with preview_container:
            ocr_result_container = ui.column().classes('w-full mt-4')
    else:
        # 清空現有內容
        ocr_result_container.clear()
    
    # 添加樣式
    ui.add_head_html('''
        <style>
            .ocr-result-card {
                margin-top: 1.5rem;
                padding: 1.5rem;
                border-radius: 8px;
                background: white;
                border: 1px solid #e0e0e0;
                box-shadow: 0 2px 8px rgba(0,0,0,0.08);
                width: 100%;
                max-width: 800px;
                margin-left: auto;
                margin-right: auto;
            }
            .ocr-result-title {
                display: flex;
                align-items: center;
                margin-bottom: 1rem;
                color: #1976D2;
                font-size: 1.1rem;
                font-weight: 500;
            }
            .ocr-result-title i {
                margin-right: 8px;
            }
            .download-btn {
                background-color: #1976D2 !important;
                color: white !important;
                padding: 8px 16px !important;
                border-radius: 4px !important;
                text-transform: none !important;
                width: auto !important;
                margin-top: 1rem;
            }
            .download-btn:hover {
                background-color: #1565C0 !important;
                box-shadow: 0 2px 4px rgba(0,0,0,0.1);
            }
            .expansion-content {
                max-height: 400px;
                overflow-y: auto;
                margin: 1rem 0;
                padding: 0.5rem;
                border: 1px solid #f0f0f0;
                border-radius: 4px;
            }
        </style>
    ''')
    
    # 創建下載按鈕的回調函數
    async def download_callback():
        nonlocal content, original_filename  # 確保可以訪問外部變數
        logger.debug("[DEBUG] 下載按鈕被點擊")
        logger.debug(f"[DEBUG] 內容類型: {type(content)}")
        logger.debug(f"[DEBUG] 內容長度: {len(content) if content else 0}")
        
        if not content:
            ui.notify("錯誤: 沒有可下載的內容", type='negative')
            return
            
        try:
            # 直接調用下載函數
            download_markdown(content, original_filename)
        except Exception as e:
            error_msg = f"下載過程中出錯: {str(e)}"
            logger.error(f"[ERROR] {error_msg}")
            ui.notify(error_msg, type='negative')
            import traceback
            traceback.print_exc()
    
    # 在結果容器中顯示內容
    with ocr_result_container:
        with ui.card().classes('ocr-result-card'):
            # 標題
            with ui.row().classes('items-center'):
                ui.icon('check_circle', color='positive').classes('text-h5')
                ui.label('OCR 處理完成').classes('text-h6 text-weight-bold q-ml-sm')
            
            # 內容預覽（可折疊）
            with ui.expansion('點擊查看完整結果', icon='unfold_more').classes('w-full q-mt-md'):
                with ui.scroll_area().classes('q-pa-sm max-h-[400px] overflow-auto'):
                    ui.markdown(content).classes('q-pa-sm')
            
            # 創建下載按鈕
            ui.button(
                '下載 Markdown 檔案',
                on_click=download_callback,
                icon='file_download',
                color='primary'
            ).classes('q-mt-md')
            
            logger.debug("[DEBUG] OCR 結果顯示完成")
    
    # 確保結果區域可見並滾動到可視區域
    ocr_result_container.visible = True
    ui.run_javascript('document.querySelector(".ocr-result-card").scrollIntoView({behavior: "smooth"})')
    ui.update(ocr_result_container)

def download_markdown(content: str, original_filename: str):
    """下載 Markdown 文件"""
    logger.debug(f"開始下載處理，原始檔名: {original_filename}")
    logger.debug(f"內容長度: {len(content) if content else 0} 字元")
    
    try:
        # 生成安全檔名
        base_name = os.path.splitext(original_filename)[0]
        safe_name = sanitize_filename(f"{base_name}_ocr.md")
        logger.debug(f"安全檔名: {safe_name}")
        
        # 確保內容是字節類型
        if isinstance(content, str):
            logger.debug("內容是字符串，進行編碼")
            content_bytes = content.encode('utf-8')
        
        logger.debug(f"內容字節長度: {len(content_bytes) if content_bytes else 0}")
        
        # 使用 JavaScript 處理下載
        js = f"""
        try {{
            // 將 Python 字節轉換為 JavaScript Uint8Array
            const content = new TextDecoder().decode(new Uint8Array({list(content_bytes)}));
            const blob = new Blob([content], {{ type: 'text/markdown;charset=utf-8' }});
            const url = URL.createObjectURL(blob);
            const a = document.createElement('a');
            a.href = url;
            a.download = '{safe_name}';
            document.body.appendChild(a);
            a.click();
            // 清理
            setTimeout(() => {{
                document.body.removeChild(a);
                URL.revokeObjectURL(url);
            }}, 100);
        }} catch (e) {{
            console.error('下載錯誤:', e);
            throw e;
        }}
        """
        ui.run_javascript(js)
        
        # 顯示成功訊息
        logger.debug(f"下載已觸發: {safe_name}")
        ui.notify(f"已開始下載: {safe_name}", type='positive')
        
    except Exception as e:
        import traceback
        error_msg = f"下載時發生錯誤: {str(e)}\n{traceback.format_exc()}"
        print(f"[ERROR] {error_msg}")
        ui.notify(f"下載時發生錯誤: {str(e)}", type='negative')
async def show_pdf_preview(file_path: Path):
    """顯示 PDF 預覽"""
    global pdf_pages, current_page
    
    # 定義回調函數
    async def prev_page():
        global current_page
        if current_page > 0:
            current_page -= 1
            await update_page()
    
    async def next_page():
        global current_page
        if current_page < pdf_pages - 1:
            current_page += 1
            await update_page()

    # 更新頁面顯示
    async def update_page():
        global current_page
        if 0 <= current_page < pdf_pages:
            # 加載頁面
            page = doc.load_page(current_page)
            # 調整縮放以獲得更好的顯示效果
            mat = fitz.Matrix(2.0, 2.0)  # 提高 DPI 以獲得更好的清晰度
            pix = page.get_pixmap(matrix=mat, alpha=False)
            
            # 將圖片保存到臨時文件
            img_path = UPLOAD_DIR / f"preview_{os.urandom(8).hex()}.png"
            pix.save(img_path)
            
            # 更新圖片
            image_container.source = str(img_path)
            page_info.text = f"PDF 頁面 {current_page + 1} / {pdf_pages}"
            
            # 更新按鈕狀態
            prev_btn.disable = current_page <= 0
            next_btn.disable = current_page >= pdf_pages - 1
                    
    try:
        # 開啟 PDF 文件
        doc = fitz.open(file_path)
        pdf_pages = len(doc)
        current_page = 0  # 重置為第一頁
        
        # 創建外層容器
        with ui.column().classes('w-full items-stretch'):
            # 創建圖片容器
            image_container = ui.image().classes('w-full max-w-4xl mx-auto border rounded')
            
            # 頁面控制按鈕
            with ui.row().classes('w-full justify-center items-center my-2') as button_row:
                # 上一頁按鈕 - 只在不是第一頁時顯示
                prev_btn = ui.button(icon='navigate_before', on_click=prev_page).props('flat dense')
                
                # 頁碼資訊
                page_info = ui.label(f"PDF 頁面 1 / {pdf_pages}").classes('mx-4')
                
                # 下一頁按鈕 - 只在不是最後一頁時顯示
                next_btn = ui.button(icon='navigate_next', on_click=next_page).props('flat dense')
            
            # 初始顯示第一頁
            await update_page()
            
    except Exception as e:
        ui.notify(f"無法預覽 PDF: {str(e)}", type='negative')

async def show_docx_preview(file_path: Path):
    """顯示 Word 文件預覽"""
    doc = DocxDocument(file_path)
    with ui.column().classes('w-full'):
        for para in doc.paragraphs:
            if para.text.strip():
                ui.markdown(f"> {para.text}")

async def show_markdown_preview(file_path: Path):
    """顯示 Markdown 文件預覽"""
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()
    ui.markdown(content)

async def show_html_preview(file_path: Path):
    """顯示 HTML 文件預覽"""
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()
    ui.html(content)

async def show_pptx_preview(file_path: Path):
    """顯示 PowerPoint 文件預覽"""
    prs = Presentation(file_path)
    with ui.column().classes('w-full'):
        for i, slide in enumerate(prs.slides):
            ui.markdown(f"### 投影片 {i+1}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    ui.markdown(f"> {shape.text}")

async def show_xlsx_preview(file_path: Path):
    """顯示 Excel 文件預覽"""
    wb = load_workbook(file_path, read_only=True)
    with ui.column().classes('w-full'):
        for sheet in wb:
            with ui.expansion(sheet.title, icon='table_chart'):
                table_data = []
                for row in sheet.iter_rows(values_only=True):
                    table_data.append(row)
                
                if table_data:
                    ui.table(
                        columns=[{'name': f'Col {i+1}', 'field': str(i), 'label': f'Col {i+1}'} for i in range(len(table_data[0]))],
                        rows=[{str(i): str(cell) if cell is not None else '' for i, cell in enumerate(row)} for row in table_data],
                        row_key='id',
                        pagination={'rowsPerPage': 10}
                    )

async def show_csv_preview(file_path: Path, num_rows=100):
    """顯示 CSV 文件預覽"""
    df = pd.read_csv(file_path)
    
    rows = [df.columns.tolist()]  # 表頭
    rows.extend(df.values.tolist())  # 資料列
    
    with ui.column().classes('w-full'):
        ui.label(f"預覽限制: {num_rows} 筆").classes('text-caption text-grey-7 q-mb-md')
        ui.table(
            columns=[{'name': str(col), 'label': str(col), 'field': str(col)} 
                    for col in df.columns],
            rows=df.head(num_rows).to_dict('records'), # 限制顯示前 100 筆資料
            pagination={'rowsPerPage': 10}
        )

# 自定義 CSS 樣式
ui.add_head_html('''
    <style>
        :root {
            --primary: #1976D2;
            --secondary: #26A69A;
            --accent: #9C27B0;
            --dark: #1D1D1D;
            --dark-page: #121212;
            --positive: #21BA45;
            --negative: #C10015;
            --info: #31CCEC;
            --warning: #F2C037;
        }
        body {
            background-color: #f5f5f5;
        }
        .custom-card {
            border-radius: 12px;
            box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
            transition: transform 0.2s, box-shadow 0.2s;
        }
        .custom-card:hover {
            transform: translateY(-2px);
            box-shadow: 0 6px 12px rgba(0, 0, 0, 0.15);
        }
        .upload-area {
            border: 2px dashed #ccc;
            border-radius: 8px;
            padding: 2rem;
            text-align: center;
            cursor: pointer;
            transition: all 0.3s;
        }
        .upload-area:hover {
            border-color: var(--primary);
            background-color: rgba(25, 118, 210, 0.05);
        }
        .file-info {
            background-color: #f8f9fa;
            border-left: 4px solid var(--primary);
            padding: 1rem;
            margin: 1rem 0;
            border-radius: 0 4px 4px 0;
        }
        .footer {
            text-align: center;
            padding: 1.5rem;
            color: #666;
            font-size: 0.9rem;
        }
        .page-controls {
            background: white;
            padding: 0.5rem;
            border-radius: 24px;
            box-shadow: 0 2px 4px rgba(0,0,0,0.1);
        }
    </style>
''')

# 創建主界面
def create_ui():
    """創建用戶界面"""
    # 重置應用程式狀態
    init_application()
    
    # 設置頁面標題和圖標
    ui.page_title("📄 文件助手 | Document Assistant")
    
    # 主容器
    with ui.column().classes('q-pa-md max-w-4xl mx-auto w-full'):
        # 頁首
        with ui.row().classes('justify-center items-center w-full q-mb-lg'):
            ui.icon('description', size='2.5rem', color='primary')
            ui.label('文件助手').classes('text-h4 text-weight-bold q-ml-sm')
        
        # 說明文字
        with ui.card().classes('w-full q-mb-md bg-blue-1'):
            with ui.row().classes('items-center'):
                ui.icon('info', color='primary')
                ui.label('支援多種文件格式預覽：PDF、Word、Excel、PPT、Markdown、HTML 等').classes('text-body2')
            with ui.row().classes('items-center'):
                ui.label(f'檔案限制：{MAX_FILE_SIZE / 1000 / 1000:.2f} MB').classes('q-ml-sm')
                ui.button('重置', on_click=init_application).classes('q-ml-sm')
        
        # 文件上傳區域
        with ui.card().classes('w-full custom-card'):
            with ui.column().classes('w-full'):
                ui.label('上傳文件').classes('text-h6 text-weight-medium q-mb-md')
                with ui.upload(
                    label='拖曳文件至此或點擊選擇',
                    on_upload=handle_upload,
                    auto_upload=True,
                    multiple=False
                ).classes('w-full') as upload:
                    with ui.column().classes('upload-area w-full'):
                        ui.icon('cloud_upload', size='3rem', color='primary')
                        ui.label('拖曳文件至此或點擊選擇').classes('q-mt-sm')
                        ui.label('(支援 PDF、Word、Excel、PPT 等格式)').classes('text-caption text-grey-7')
        
        # 預覽區域
        global preview_container
        preview_container = ui.column().classes('w-full q-mt-lg')
        
        # 頁尾
        with ui.row().classes('w-full justify-center q-mt-xl'):
            ui.label('© 2025 Document Assistant').classes('text-caption text-grey-7')
    
    # 添加全局加載狀態
    global loading
    loading = ui.linear_progress(show_value=False, size='2px', color='primary')
    loading.visible = False

# 啟動應用
if __name__ in ["__main__", "__mp_main__"]:
    app.add_static_files('/temp_uploads', 'temp_uploads')
    create_ui()
    ui.run(title="Document Assistant", port=8080, reload=False, show=False)
