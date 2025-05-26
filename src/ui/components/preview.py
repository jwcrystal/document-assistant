"""
文件預覽組件模組

此模組提供不同類型文件的預覽功能。
"""
import fitz  # PyMuPDF
import pandas as pd
from pathlib import Path
from typing import Optional
from nicegui import ui

from src.config import settings

class FilePreview:
    """文件預覽基類"""
    
    def __init__(self, file_path: Path):
        self.file_path = file_path
    
    async def show(self):
        """顯示文件預覽"""
        raise NotImplementedError("子類必須實現此方法")


class PDFPreview(FilePreview):
    """PDF 文件預覽"""
    
    async def show(self):
        """顯示 PDF 預覽"""
        try:
            # 使用 PyMuPDF 提取 PDF 第一頁作為預覽
            doc = fitz.open(self.file_path)
            
            # 顯示頁面導航
            with ui.row().classes('w-full justify-center'):
                page_slider = ui.slider(min=1, max=len(doc), value=1, step=1).props('label-slot')
                with page_slider.add_slot('label'):
                    ui.label('頁面')
                ui.label().bind_text_from(page_slider, 'value')
            
            # 顯示當前頁面
            page_container = ui.column().classes('w-full items-center')
            
            async def update_page(e):
                # 從事件參數中獲取值
                value = e.args if isinstance(e.args, (int, float)) else getattr(e, 'value', 1)
                page_num = int(value) - 1
                if 0 <= page_num < len(doc):
                    page = doc.load_page(page_num)
                    pix = page.get_pixmap()
                    img_path = settings.UPLOAD_DIR / f"preview_{self.file_path.stem}_{page_num}.png"
                    pix.save(img_path)
                    
                    page_container.clear()
                    with page_container:
                        ui.image(str(img_path)).classes('max-w-full border')
            
            page_slider.on('update:model-value', update_page)
            
            # 觸發初始頁面加載
            await update_page(type('obj', (), {'args': 1}))
            
        except Exception as e:
            ui.notify(f"預覽 PDF 時出錯: {str(e)}", type='negative')
            ui.label(f"無法預覽 PDF: {str(e)}")


class DocxPreview(FilePreview):
    """Word 文件預覽"""
    
    async def show(self):
        """顯示 Word 文件預覽"""
        try:
            from docx import Document
            doc = Document(self.file_path)
            
            # 提取文本內容
            content = []
            for para in doc.paragraphs:
                if para.text.strip():
                    content.append(para.text)
            
            # 顯示預覽
            with ui.column().classes('w-full border p-4'):
                ui.label("預覽 (純文本):").classes('text-lg font-bold')
                ui.separator()
                ui.markdown('\n\n'.join(content[:20]))  # 只顯示前20段
                if len(content) > 20:
                    ui.label("... (僅顯示部分內容)")
                    
        except Exception as e:
            ui.notify(f"預覽 Word 文檔時出錯: {str(e)}", type='negative')
            ui.label(f"無法預覽 Word 文檔: {str(e)}")


class MarkdownPreview(FilePreview):
    """Markdown 文件預覽"""
    
    async def show(self):
        """顯示 Markdown 預覽"""
        try:
            with open(self.file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            with ui.column().classes('w-full'):
                ui.markdown(content)
                
        except Exception as e:
            ui.notify(f"預覽 Markdown 時出錯: {str(e)}", type='negative')
            ui.label(f"無法預覽 Markdown 文件: {str(e)}")


class HTMLPreview(FilePreview):
    """HTML 文件預覽"""
    
    async def show(self):
        """顯示 HTML 預覽"""
        try:
            with open(self.file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            with ui.column().classes('w-full h-96'):
                ui.html(content).classes('w-full h-full border')
                
        except Exception as e:
            ui.notify(f"預覽 HTML 時出錯: {str(e)}", type='negative')
            ui.label(f"無法預覽 HTML 文件: {str(e)}")


class PPTXPreview(FilePreview):
    """PowerPoint 文件預覽"""
    
    async def show(self):
        """顯示 PowerPoint 預覽"""
        try:
            from pptx import Presentation
            prs = Presentation(self.file_path)
            
            with ui.column().classes('w-full'):
                ui.label(f"幻燈片總數: {len(prs.slides)}")
                
                # 顯示每頁的縮略圖和備註
                for i, slide in enumerate(prs.slides[:5]):  # 只顯示前5頁
                    with ui.expansion(f"幻燈片 {i+1}", icon='slideshow').classes('w-full'):
                        with ui.column().classes('w-full border p-2'):
                            # 顯示幻燈片備註
                            notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else "(無備註)"
                            ui.label(f"備註: {notes[:200]}" + ("..." if len(notes) > 200 else ""))
                
                if len(prs.slides) > 5:
                    ui.label(f"... 還有 {len(prs.slides)-5} 頁未顯示")
                    
        except Exception as e:
            ui.notify(f"預覽 PowerPoint 時出錯: {str(e)}", type='negative')
            ui.label(f"無法預覽 PowerPoint 文件: {str(e)}")


class ExcelPreview(FilePreview):
    """Excel 文件預覽"""
    
    async def show(self, num_rows: int = 100):
        """顯示 Excel 預覽"""
        try:
            xl = pd.ExcelFile(self.file_path)
            
            with ui.column().classes('w-full'):
                # 工作表選擇器
                with ui.row().classes('w-full items-center'):
                    sheet_select = ui.select(
                        label='工作表',
                        options=xl.sheet_names,
                        value=xl.sheet_names[0],
                    ).classes('w-64')
                    
                    # 添加展開/收合按鈕
                    expand_btn = ui.button('展開表格', on_click=None).props('flat')
                    
                # 創建一個可收合的容器
                with ui.card().classes('w-full') as card:
                    # 添加一個過渡效果
                    card.classes('transition-all duration-300')
                    
                    # 初始狀態為收合
                    is_expanded = False
                    
                    # 表格容器
                    table_container = ui.column().classes('w-full')
                    
                    def update_table(sheet_name: str):
                        try:
                            # 讀取選中的工作表
                            df = xl.parse(sheet_name)
                            
                            # 清空容器
                            table_container.clear()
                            
                            # 顯示數據總行數
                            with table_container:
                                ui.label(f'總行數: {len(df)}')
                                
                                # 創建分頁表格
                                grid = ui.aggrid.from_pandas(df.head(num_rows))  # 限制最大顯示100行
                                grid.classes('w-full')
                                
                                # 添加分頁控制
                                grid.options['pagination'] = True
                                grid.options['paginationPageSize'] = 10
                                grid.update()
                                
                                if len(df) > num_rows:
                                    ui.label(f"... (僅顯示前{num_rows}行)")
                                    
                            return True
                        except Exception as e:
                            ui.notify(f'載入工作表 {sheet_name} 時出錯: {str(e)}', type='negative')
                            return False
                    
                    # 切換展開/收合狀態的函數
                    def toggle_expand():
                        nonlocal is_expanded
                        is_expanded = not is_expanded
                        if is_expanded:
                            card.classes(remove='max-h-12', add='max-h-screen')
                            expand_btn.text = '收合表格'
                            # 加載數據
                            ui.timer(0.1, lambda: update_table(sheet_select.value), once=True)
                        else:
                            card.classes(remove='max-h-screen', add='max-h-12')
                            expand_btn.text = '展開表格'
                    
                    # 綁定按鈕點擊事件
                    expand_btn.on_click(toggle_expand)
                    
                    # 綁定工作表選擇變化事件
                    def on_sheet_change(e):
                        # 從事件參數中獲取工作表名稱
                        if isinstance(e.args, dict) and 'label' in e.args:
                            update_table(e.args['label'])
                        else:
                            update_table(e.args)
                    
                    sheet_select.on('update:model-value', on_sheet_change)
                    
                    # 初始狀態為收合
                    card.classes('max-h-12 overflow-hidden transition-all duration-300')
                    
                    # 初始加載第一個工作表的數據
                    ui.timer(0.1, lambda: update_table(xl.sheet_names[0]), once=True)
                    
        except Exception as e:
            ui.notify(f"預覽 Excel 時出錯: {str(e)}", type='negative')
            ui.label(f"無法預覽 Excel 文件: {str(e)}")


class CSVPreview(FilePreview):
    """CSV 文件預覽"""
    
    async def show(self, num_rows: int = 100):
        """顯示 CSV 預覽
        
        Args:
            num_rows: 要顯示的行數
        """
        try:
            df = pd.read_csv(self.file_path, nrows=num_rows)
            
            with ui.column().classes('w-full'):
                # 創建表格
                columns = [{'name': col, 'label': col, 'field': col} for col in df.columns]
                rows = df.to_dict('records')
                
                # 創建表格並設置列和行
                table = ui.table(
                    columns=columns,
                    rows=rows,
                    row_key='id',
                    pagination={'rowsPerPage': 10}
                ).classes('w-full')
                
                if len(df) >= num_rows:
                    ui.label(f"... (僅顯示前{num_rows}行)")
                    
        except Exception as e:
            ui.notify(f"預覽 CSV 時出錯: {str(e)}", type='negative')
            ui.label(f"無法預覽 CSV 文件: {str(e)}")


# 文件類型到預覽類的映射
PREVIEW_CLASSES = {
    'application/pdf': PDFPreview,
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': DocxPreview,
    'application/msword': DocxPreview,
    'text/markdown': MarkdownPreview,
    'text/html': HTMLPreview,
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': PPTXPreview,
    'application/vnd.ms-powerpoint': PPTXPreview,
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': ExcelPreview,
    'application/vnd.ms-excel': ExcelPreview,
    'text/csv': CSVPreview,
}


def get_preview_handler(file_path: Path, file_type: str) -> Optional[FilePreview]:
    """
    根據文件類型獲取對應的預覽處理器
    
    Args:
        file_path: 文件路徑
        file_type: 文件 MIME 類型
        
    Returns:
        FilePreview 實例，如果不支援則返回 None
    """
    preview_class = PREVIEW_CLASSES.get(file_type)
    if not preview_class:
        return None
    return preview_class(file_path)
